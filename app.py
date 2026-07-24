import streamlit as st
import chromadb
from sentence_transformers import SentenceTransformer
from groq import Groq
from pypdf import PdfReader
from docx import Document
import openpyxl
from pptx import Presentation
import io

# ---- PAGE SETUP ----
st.set_page_config(page_title="Company Knowledge Assistant", page_icon="📄")
st.title("📄 Company Knowledge Assistant")
st.write("Upload company documents, then ask questions about them.")

# ---- LOAD MODELS (cached so it only loads once) ----
@st.cache_resource
def load_embedder():
    return SentenceTransformer('all-MiniLM-L6-v2')

@st.cache_resource
def get_groq_client():

    return Groq(api_key=st.secrets["GROQ_API_KEY"])

embedder = load_embedder()
client = get_groq_client()

# ---- HELPER FUNCTIONS TO READ DIFFERENT FILE TYPES (from uploaded file objects) ----
def read_txt(file_obj):
    return file_obj.read().decode("utf-8")

def read_pdf(file_obj):
    reader = PdfReader(file_obj)
    text = ""
    for page in reader.pages:
        text += page.extract_text() or ""
    return text

def read_docx(file_obj):
    doc = Document(file_obj)
    return "\n".join([para.text for para in doc.paragraphs])

def read_xlsx(file_obj):
    wb = openpyxl.load_workbook(file_obj, data_only=True)
    text = ""
    for sheet in wb.worksheets:
        text += f"\n--- Sheet: {sheet.title} ---\n"
        for row in sheet.iter_rows(values_only=True):
            row_text = " | ".join(str(cell) for cell in row if cell is not None)
            if row_text.strip():
                text += row_text + "\n"
    return text

def read_pptx(file_obj):
    prs = Presentation(file_obj)
    text = ""
    for i, slide in enumerate(prs.slides, start=1):
        text += f"\n--- Slide {i} ---\n"
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = "".join(run.text for run in para.runs)
                    if line.strip():
                        text += line + "\n"
    return text

def extract_text(uploaded_file):
    name = uploaded_file.name
    file_bytes = io.BytesIO(uploaded_file.read())
    if name.endswith(".txt"):
        return read_txt(file_bytes)
    elif name.endswith(".pdf"):
        return read_pdf(file_bytes)
    elif name.endswith(".docx"):
        return read_docx(file_bytes)
    elif name.endswith(".xlsx"):
        return read_xlsx(file_bytes)
    elif name.endswith(".pptx"):
        return read_pptx(file_bytes)
    else:
        return ""

# ---- SESSION STATE SETUP ----
if "messages" not in st.session_state:
    st.session_state.messages = []
if "collection" not in st.session_state:
    st.session_state.collection = None
if "loaded_files" not in st.session_state:
    st.session_state.loaded_files = []

# ---- FILE UPLOAD UI ----
st.subheader("1. Upload your documents")
uploaded_files = st.file_uploader(
    "Upload PDF, Word, Excel, PowerPoint, or text files",
    type=["txt", "pdf", "docx", "xlsx", "pptx"],
    accept_multiple_files=True,
)

if st.button("Process documents"):
    if uploaded_files:
        with st.spinner("Reading and indexing documents..."):
            documents = []
            filenames = []

            for uploaded_file in uploaded_files:
                try:
                    text = extract_text(uploaded_file)
                    if text.strip():
                        documents.append(text)
                        filenames.append(uploaded_file.name)
                    else:
                        st.warning(f"No text extracted from {uploaded_file.name} (might be a scanned/image file)")
                except Exception as e:
                    st.error(f"Error reading {uploaded_file.name}: {e}")

            if documents:
                chroma_client = chromadb.Client()
                # Use get_or_create to avoid "already exists" errors on rerun
                collection = chroma_client.get_or_create_collection(name="session_docs")
                embeddings = embedder.encode(documents).tolist()
                collection.add(documents=documents, embeddings=embeddings, ids=filenames)

                st.session_state.collection = collection
                st.session_state.loaded_files = filenames
                st.session_state.messages = []  # reset chat when new documents are loaded

        st.success(f"Loaded {len(st.session_state.loaded_files)} documents: {', '.join(st.session_state.loaded_files)}")
    else:
        st.warning("Please upload at least one file first.")

# ---- CHAT SECTION (only shown once documents are processed) ----
if st.session_state.collection is not None:
    st.subheader("2. Ask questions")

    for message in st.session_state.messages:
        with st.chat_message(message["role"]):
            st.write(message["content"])

    question = st.chat_input("Type your question here...")

    if question:
        st.session_state.messages.append({"role": "user", "content": question})
        with st.chat_message("user"):
            st.write(question)

        with st.chat_message("assistant"):
            with st.spinner("Searching documents and generating answer..."):
                question_embedding = embedder.encode([question]).tolist()

                results = st.session_state.collection.query(
                    query_embeddings=question_embedding,
                    n_results=2
                )

                retrieved_chunks = results['documents'][0]
                retrieved_files = results['ids'][0]

                context = "\n\n".join(retrieved_chunks)

                prompt = f"""Answer the question using ONLY the information below. If the answer isn't in the information provided, say so.

INFORMATION:
{context}

QUESTION: {question}
"""

                response = client.chat.completions.create(
                    model="llama-3.3-70b-versatile",
                    max_tokens=500,
                    messages=[{"role": "user", "content": prompt}]
                )

                answer = response.choices[0].message.content
                st.write(answer)

                with st.expander("See which documents were used"):
                    st.write(retrieved_files)

        st.session_state.messages.append({"role": "assistant", "content": answer})
else:
    st.info("Upload documents and click 'Process documents' to start asking questions.")